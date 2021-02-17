# -*- coding: utf-8 -*-
import re
import os
import logging
import shutil
from collections import Counter
import random
from multiprocessing import Pool as ProcessPool
from multiprocessing.dummy import Pool as ThreadPool
from queue import Queue as ThreadQueue

import jieba
from jieba import posseg

from .cmd_tools import tika_convert, antiword_convert
from . import _config
from .data_helps import corpus_cut

__all__ = (
    "Files2Text",
    "CorpusClean",
    "CorpusHandler"
)

jieba.setLogLevel(logging.INFO)
jieba.enable_parallel()
jieba_inst = posseg.POSTokenizer(jieba.Tokenizer())


class ReCache:
    N_RE = re.compile(r'\n')
    MORE_SPACE_T_RE = re.compile(r'[\s\t]+')
    CN_RE = re.compile(r'^[\u4e00-\u9fa5]+$')


class BaseCorpus(object):

    def __init__(self, logger=None):
        if logger is None:
            logging.basicConfig(level=logging.DEBUG,
                                format=_config.LOG_FORMAT,
                                datefmt=_config.LOG_DATE_FORMAT)
            logger = logging.getLogger(BaseCorpus.__class__.__name__)
        self.logger = logger


class Files2Text(BaseCorpus):
    """
    各文件转换成文本语料
    文件格式支持：doc、docx、txt、html、xml、pdf、ppt、xls、xlsx、pptx、png、jpg
    """

    _POOL_UPPER_LIMIT = 1000

    def __init__(self, re_cache_cls=None, **kwargs):
        super(Files2Text, self).__init__(**kwargs)
        if re_cache_cls is None:
            re_cache_cls = ReCache
        self._re_cache = re_cache_cls()
        self._cut_words = False
        self._pool_size = os.cpu_count()

    @staticmethod
    def _doc_to_text(filename):
        if filename.endswith(".doc"):
            try:
                return antiword_convert(filename)
            except:
                return tika_convert(filename)
        else:
            import docx
            text_content = docx.Document(filename)
            return "".join(row.text.strip() for row in text_content.paragraphs if row.text.strip())

    def _pptx_txt_from_shape(self, shape):
        from pptx.shapes.group import GroupShape
        ret = ""
        if type(shape) == GroupShape:
            for sshape in shape.shapes:
                self._pptx_txt_from_shape(sshape)
        elif shape.has_text_frame:
            ret += "".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
        elif shape.has_table:
            for row in shape.table.rows:
                bfirst = True
                for cell in (_cell for _cell in row.cells if _cell.text_frame):
                    content = cell.text.strip()
                    if content:
                        if bfirst:
                            ret += content
                            bfirst = False
                        else:
                            ret += content
        return ret

    def _pptx_to_text(self, filename):
        import pptx
        ppt = pptx.Presentation(filename)
        content = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                content += self._pptx_txt_from_shape(shape)
        return content

    @staticmethod
    def _image_to_text(filename):
        from .aip import AipOcr
        app = AipOcr(_config.IMAGES_DISCERN_CONF["app_id"],
                     _config.IMAGES_DISCERN_CONF["api_key"],
                     _config.IMAGES_DISCERN_CONF["secret_key"])
        with open(filename, 'rb') as fp:
            result = app.basicGeneral(fp)
            words_result = result.get("words_result", None)
            return "".join(row["words"].strip() for row in words_result if row["words"].strip())

    def _file_discern(self, filename):
        if filename.lower().endswith((".doc", ".docx")):
            return self._doc_to_text(filename)
        elif filename.lower().endswith((".txt", ".html", ".xml", ".pdf", ".ppt", ".xls", ".xlsx")):
            return tika_convert(filename)
        elif filename.lower().endswith(".pptx"):
            return self._pptx_to_text(filename)
        elif filename.lower().endswith((".png", "jpg")):
            return self._image_to_text(filename)
        else:
            self.logger.warning("The %s cannot be recognized....", filename)

    def _task_callback(self, f, sample_path, filename, file_content_split, clean_func=None, filename_label_list=None):
        try:
            file_path = os.path.join(sample_path, filename)
            sample_content = self._file_discern(file_path)
            if not sample_content:
                return filename
            sample_content = self._re_cache.N_RE.sub("", sample_content)
            sample_content = self._re_cache.MORE_SPACE_T_RE.sub(" ", sample_content)
            if clean_func:
                sample_content = clean_func(sample_content)
                if not sample_content:
                    return filename
            content = filename + file_content_split + sample_content + "\n"
            if filename_label_list is not None:
                label = str(dict(filename_label_list)[filename]).strip()
                content = label + file_content_split + content
            self.logger.debug(filename)
        except:
            import traceback
            self.logger.error(traceback.format_exc())
        else:
            f.write(content)

    def _gen_callback_args(self, sample_path, max_files_num, f, file_content_split, filename_label_list):
        for index, item in enumerate(filename_label_list):
            if index > max_files_num:
                break
            filename = item[0]
            if self._cut_words and filename_label_list is not None:
                yield f, sample_path, filename, file_content_split, self._clean_func, filename_label_list
            elif self._cut_words and filename_label_list is None:
                yield f, sample_path, filename, file_content_split, self._clean_func
            else:
                yield f, sample_path, filename, file_content_split

    def _validate_data(self, pool_size, sample_path):
        if not sample_path:
            raise ValueError("a sample_path is required!")
        if sample_path and not os.path.isdir(sample_path):
            raise ValueError("sample_path must be a path!")
        if pool_size and pool_size < self._POOL_UPPER_LIMIT:
            self._pool_size = pool_size

    def start(self,
              sample_path,
              output_file=_config.DEFAULT_OUTPUT_FILE,
              max_files_num=5000,
              file_content_split=_config.DEFAULT_FILE_CONTENT_SPLIT,
              pool_size=None,
              filename_label_list=None):

        self._validate_data(pool_size, sample_path)
        self.logger.info("count samples: %s", len(filename_label_list))
        self.logger.info("discern sample and generate corpus ....")
        filter_sample = list()
        with open(output_file, "a", encoding="utf-8") as f:
            for params in self._gen_callback_args(sample_path, max_files_num, f, file_content_split, filename_label_list):
                filename = self._task_callback(*params)
                if filename:
                    filter_sample.append(filename)
        return filter_sample

    def __call__(self, clean_func, *args, **kwargs):
        self._cut_words = True
        self._clean_func = clean_func
        return self


class CorpusClean(object):
    """
    语料清洗
    只保留中文，且只关注名词、动词、形容词、数量词、量词、代词
    """

    def __init__(self, re_cache_cls=None):
        if re_cache_cls is None:
            re_cache_cls = ReCache
        self._base_path = os.path.dirname(os.path.abspath(__file__))
        self._re_cache = re_cache_cls()
        self._jieba_inst = None
        self._stopwords = None
        self._thread_queue = None

    @staticmethod
    def _cut_words(sample_content):
        return jieba_inst.cut(sample_content)

    @staticmethod
    def _pos_tag(words):
        for w in words:
            pos = w.flag[0]
            if pos not in _config.RETAIN_POS.keys():
                continue
            yield w.word

    def _retain_cn(self, words):
        for w in words:
            if not self._re_cache.CN_RE.search(w):
                continue
            yield w

    def _remove_stopwords(self, words):
        for w in words:
            if w in self._stopwords:
                continue
            yield w

    @staticmethod
    def _remove_one_word(words):
        for w in words:
            if len(w) < 2:
                continue
            yield w

    @staticmethod
    def _merge_words(words):
        return " ".join(w for w in words)

    def _load_stopwords(self):
        stopwords_file = os.path.join(self._base_path, "stop_words/stop_words.txt")
        stopwords = set()
        with open(stopwords_file, "r", encoding="utf-8") as f:
            for w in f.readlines():
                stopwords.add(w.strip())
        return stopwords

    def load_resource(self):
        self._stopwords = self._load_stopwords()

    def clean_handler(self, sample_content):
        words = self._cut_words(sample_content)
        words = self._pos_tag(words)
        words = self._retain_cn(words)
        words = self._remove_stopwords(words)
        words = self._remove_one_word(words)
        sample_content = self._merge_words(words)
        return sample_content

    @staticmethod
    def _validate_data(sample_content):
        if not sample_content:
            raise ValueError("a sample_content is required!")
        if sample_content and not isinstance(sample_content, list):
            raise ValueError("sample_content must be a list type")

    def _process_task(self, *params):
        return [(index, self.clean_handler(sample)) for index, sample in zip(*params)]

    def _thread_task(self, params):
        for index, sample in params:
            self._thread_queue.put((index, self.clean_handler(sample)))

    def _start_process(self, more_sample_content):
        cpu = os.cpu_count()
        pool = ProcessPool(cpu)
        index_batch, sample_batch = corpus_cut(more_sample_content, cpu)
        result = [pool.apply_async(self._process_task, args=(index_batch[i], sample_batch[i])) for i in range(cpu)]
        pool.close()
        pool.join()
        result = map(lambda item: item[1], sorted([item for res in result for item in res.get()],
                                                  key=lambda item: item[0]))
        return list(result)

    def _start_thread(self, more_sample_content):
        pool = ThreadPool(len(more_sample_content))
        self._thread_queue = ThreadQueue()
        sample_data = [(index, sample) for index, sample in enumerate(more_sample_content)]
        pool.apply_async(self._thread_task, args=(sample_data, ))
        pool.close()
        pool.join()
        result = list()
        while not self._thread_queue.empty():
            result.append(self._thread_queue.get())
            self._thread_queue.task_done()
        result = map(lambda item: item[1], sorted(result, key=lambda item: item[0]))
        return list(result)

    def start(self, more_sample_content: list, mul_task):
        self._validate_data(more_sample_content)
        if mul_task == "process_pool":
            return self._start_process(more_sample_content)
        elif mul_task == "thread_pool":
            return self._start_thread(more_sample_content)
        else:
            return [self.clean_handler(item) for item in more_sample_content]


class CorpusResult(object):

    __slots__ = ("train_corpus", "validate_corpus", "train_label", "validate_label")


class CorpusHandler(BaseCorpus):
    """语料分析"""

    def __init__(self, logger=None, re_cache_cls=None):
        super(CorpusHandler, self).__init__(logger)
        if re_cache_cls is None:
            re_cache_cls = ReCache
        self._file2text_inst = Files2Text(logger=logger, re_cache_cls=re_cache_cls)
        self._corpus_clean_inst = CorpusClean(re_cache_cls=re_cache_cls)
        self._filter_filename = list()
        self._corpus_result = CorpusResult()

    def cp_sample(self, sample_path, tmp_path, filename_label_list):
        """复制本地文件到临时目录，安全考虑"""
        for filename_by_path, label in filename_label_list:
            label_path, filename = os.path.split(filename_by_path)
            local_file = os.path.join(os.path.join(sample_path, label_path), filename)
            label_tmp_path = os.path.join(tmp_path, label_path)
            if not os.path.exists(label_tmp_path):
                os.makedirs(label_tmp_path, exist_ok=True)
            tmp_file = os.path.join(label_tmp_path, filename)
            try:
                shutil.copyfile(local_file, tmp_file)
            except:
                import traceback
                self.logger.error(traceback.format_exc())

    def get_corpus(self, corpus_data: list, mul_task):
        """
        获取预测语料数据，为模型预测做准备
        """
        return self._corpus_clean_inst.start(corpus_data, mul_task)

    def get_train_and_validate_corpus_by_extract(self, file_path_result, sample_path, filename_label_list,
                                                 validate_percentage=.3, label_min_sample_num=5):
        """
        获取训练语料数据，为建模做准备
        :param file_path_result: 文件路径处理结果
        :param sample_path: 真实样本路径
        :param filename_label_list: 对接后台的文件名和类别数据
        :param validate_percentage: 验证集比例大小，默认为0.3
        :param label_min_sample_num: 单个类别最少样本数
        :return: 返回建模所需预料数据
        """
        self.gen_corpus(sample_path, file_path_result.tmp_sample_path, filename_label_list,
                        file_path_result.tmp_corpus_file)
        self._filter_sample(file_path_result.tmp_corpus_file, filename_label_list, extract_rules=True)
        corpus_result = self.split_data(file_path_result.tmp_corpus_file, validate_percentage, label_min_sample_num)
        return corpus_result

    def get_train_and_validate_corpus_by_add(self, file_path_result, sample_path, filename_label_list,
                                             validate_percentage=.3, label_min_sample_num=5):
        from .data_helps import add_corpus_content_and_filename
        filter_sample = self.gen_corpus(sample_path, file_path_result.tmp_sample_path, filename_label_list,
                                        file_path_result.tmp_corpus_file)
        add_corpus_content_and_filename(file_path_result.tmp_corpus_file, file_path_result.corpus_data_file)
        self._filter_sample(file_path_result.corpus_data_file, extract_rules=False, filter_sample=filter_sample)
        corpus_result = self.split_data(file_path_result.corpus_data_file, validate_percentage, label_min_sample_num)
        return corpus_result

    def get_train_and_validate_corpus_by_del(self, corpus_data_file, filename_label_list, validate_percentage=.3,
                                             label_min_sample_num=5):
        from .data_helps import del_corpus_content_by_filename
        del_corpus_content_by_filename(corpus_data_file, filename_label_list)
        corpus_result = self.split_data(corpus_data_file, validate_percentage, label_min_sample_num)
        return corpus_result

    def get_train_and_validate_corpus_by_semi(self, file_path_result, sample_path, filename_label_list):
        from .data_helps import read_tmp_corpus
        self.gen_corpus(sample_path, file_path_result.tmp_sample_path, filename_label_list,
                        file_path_result.tmp_corpus_file)
        semi_corpus_result = read_tmp_corpus(file_path_result.tmp_corpus_file)
        return semi_corpus_result

    def gen_corpus(self, sample_path, tmp_sample_path, filename_label_list, tmp_corpus_file):
        self.cp_sample(sample_path, tmp_sample_path, filename_label_list)
        handler = self._file2text_inst(self._corpus_clean_inst.clean_handler)
        filter_sample = handler.start(tmp_sample_path, output_file=tmp_corpus_file,
                                      filename_label_list=filename_label_list)
        return filter_sample

    def split_data(self, tmp_corpus_file, validate_percentage, label_min_sample_num, seed=51):
        """分割语料库为训练集和验证集"""
        from .data_helps import split_corpus, split_label_and_content, min_label_count
        while True:
            train_data, validate_data = split_corpus(tmp_corpus_file, validate_percentage=validate_percentage,
                                                     seed=seed)
            train_label, train_content = split_label_and_content(train_data)
            validate_label, validate_content = split_label_and_content(validate_data)
            min_label_key, min_label_value = min_label_count(train_label + validate_label)
            if min_label_value < label_min_sample_num:
                raise ValueError("after corpus analysis, the number of category %s is less than %s" % (
                    min_label_key, label_min_sample_num))
            train_label_count = Counter(train_label).values()
            self.logger.debug("seed: %s, train_label_count: %s", seed, train_label_count)
            if min(train_label_count) < label_min_sample_num:
                seed = random.choice(range(40, 100))
                continue
            self._corpus_result.train_corpus = train_content
            self._corpus_result.validate_corpus = validate_content
            self._corpus_result.train_label = train_label
            self._corpus_result.validate_label = validate_label
            return self._corpus_result

    def _filter_sample(self, corpus_file, filename_label_list=None, extract_rules=True, filter_sample: list=None):
        from .data_helps import get_filter_sample
        filter_filename = get_filter_sample(corpus_file, filename_label_list, extract_rules)
        if filter_sample:
            filter_filename = filter_filename + filter_sample
        self._filter_filename = filter_filename

    @property
    def filter_filename(self):
        return self._filter_filename

    @property
    def corpus_clean_tool(self):
        return self._corpus_clean_inst
